from datetime import date, timedelta
import io

import pandas as pd
import plotly.express as px
import streamlit as st
import yfinance as yf
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

st.set_page_config(page_title="Portfolio Analyser", layout="wide")

REGION_ORDER = [
    "North America",
    "Europe (excluding UK)",
    "UK",
    "Japan",
    "Rest of World",
    "Unclassified",
]

ASSET_CLASS_ORDER = [
    "Equity",
    "Bonds",
    "Alternatives",
    "Cash",
    "Unclassified",
]

COLUMN_ALIASES = {
    "identifier": ["ticker", "isin", "sedol", "security code", "instrument"],
    "description": ["description", "holding", "security name", "name"],
    "gbp_value": ["gbp value", "market value", "value", "market value gbp", "gbp_market_value"],
}

MANUAL_EXCEPTIONS = {
    "CASHGBP": {
        "matched_ticker": None,
        "asset_class": "Cash",
        "region": "UK",
        "match_status": "Manual exception",
    },
    "GBP CASH": {
        "matched_ticker": None,
        "asset_class": "Cash",
        "region": "UK",
        "match_status": "Manual exception",
    },
}


def find_matching_column(columns, aliases):
    lower_map = {c.lower().strip(): c for c in columns}
    for alias in aliases:
        if alias in lower_map:
            return lower_map[alias]
    return None


def detect_identifier_type(raw):
    raw = str(raw).strip()
    if len(raw) == 12 and raw[:2].isalpha():
        return "ISIN"
    if len(raw) == 7 and raw.isalnum():
        return "SEDOL"
    return "Ticker"


def parse_uploaded_holdings(uploaded_file):
    df = pd.read_csv(uploaded_file)

    identifier_col = find_matching_column(df.columns, COLUMN_ALIASES["identifier"])
    description_col = find_matching_column(df.columns, COLUMN_ALIASES["description"])
    value_col = find_matching_column(df.columns, COLUMN_ALIASES["gbp_value"])

    missing = []
    if not identifier_col:
        missing.append("identifier")
    if not description_col:
        missing.append("description")
    if not value_col:
        missing.append("gbp_value")

    if missing:
        raise ValueError(
            "Missing required columns: " + ", ".join(missing) + ". "
            "Your file needs a holding identifier, description, and GBP value column."
        )

    parsed = df[[identifier_col, description_col, value_col]].copy()
    parsed.columns = ["identifier", "description", "gbp_value"]

    parsed["identifier"] = parsed["identifier"].astype(str).str.strip()
    parsed["description"] = parsed["description"].astype(str).str.strip()
    parsed["gbp_value"] = (
        parsed["gbp_value"]
        .astype(str)
        .str.replace(",", "", regex=False)
        .str.replace("£", "", regex=False)
        .astype(float)
    )
    parsed["identifier_type"] = parsed["identifier"].apply(detect_identifier_type)

    return parsed


@st.cache_data(show_spinner=False)
def get_ticker_metadata(ticker):
    try:
        info = yf.Ticker(ticker).info
        if not isinstance(info, dict):
            return {}
        return info
    except Exception:
        return {}


@st.cache_data(show_spinner=False)
def get_price_history(ticker: str, years: int):
    if ticker is None or str(ticker).strip() == "":
        return pd.Series(dtype=float)

    end_date = date.today()
    start_date = end_date - timedelta(days=365 * years + 10)

    try:
        data = yf.download(
            ticker,
            start=start_date,
            end=end_date,
            auto_adjust=True,
            progress=False,
        )
    except Exception:
        return pd.Series(dtype=float)

    if data.empty:
        return pd.Series(dtype=float)

    close_series = data["Close"].copy()
    if hasattr(close_series, "squeeze"):
        close_series = close_series.squeeze()

    return close_series


def map_country_to_region(country):
    if not country:
        return "Unclassified"

    country = str(country).strip().lower()

    north_america = {
        "united states", "usa", "us", "canada", "mexico"
    }
    uk = {
        "united kingdom", "uk", "great britain", "england", "scotland", "wales", "northern ireland"
    }
    europe_ex_uk = {
        "france", "germany", "switzerland", "netherlands", "sweden", "norway",
        "denmark", "finland", "italy", "spain", "belgium", "austria", "ireland",
        "portugal", "luxembourg", "poland"
    }
    japan = {"japan"}

    if country in north_america:
        return "North America"
    if country in uk:
        return "UK"
    if country in europe_ex_uk:
        return "Europe (excluding UK)"
    if country in japan:
        return "Japan"

    return "Rest of World"


def map_quote_type_to_asset_class(quote_type, sector, long_name):
    qt = str(quote_type).strip().upper() if quote_type else ""
    sector = str(sector).strip().lower() if sector else ""
    long_name = str(long_name).strip().lower() if long_name else ""

    if qt in {"EQUITY", "ETF", "MUTUALFUND"}:
        if "bond" in long_name or sector == "fixed income":
            return "Bonds"
        return "Equity"

    if qt in {"BOND"}:
        return "Bonds"

    if "bond" in long_name:
        return "Bonds"
    if "gold" in long_name or "commodity" in long_name:
        return "Alternatives"

    return "Unclassified"


def resolve_single_holding(identifier, description):
    identifier = str(identifier).strip()
    description = str(description).strip()

    if identifier in MANUAL_EXCEPTIONS:
        manual = MANUAL_EXCEPTIONS[identifier]
        return {
            "matched_ticker": manual["matched_ticker"],
            "asset_class": manual["asset_class"],
            "region": manual["region"],
            "match_status": manual["match_status"],
            "metadata_country": None,
            "metadata_quote_type": None,
        }

    identifier_type = detect_identifier_type(identifier)

    if identifier_type != "Ticker":
        return {
            "matched_ticker": None,
            "asset_class": "Unclassified",
            "region": "Unclassified",
            "match_status": f"{identifier_type} not resolved",
            "metadata_country": None,
            "metadata_quote_type": None,
        }

    info = get_ticker_metadata(identifier)

    if not info:
        return {
            "matched_ticker": identifier,
            "asset_class": "Unclassified",
            "region": "Unclassified",
            "match_status": "Ticker found, metadata missing",
            "metadata_country": None,
            "metadata_quote_type": None,
        }

    country = info.get("country")
    quote_type = info.get("quoteType")
    sector = info.get("sector")
    long_name = info.get("longName") or info.get("shortName") or description

    region = map_country_to_region(country)
    asset_class = map_quote_type_to_asset_class(quote_type, sector, long_name)

    return {
        "matched_ticker": identifier,
        "asset_class": asset_class,
        "region": region,
        "match_status": "Yahoo Finance metadata",
        "metadata_country": country,
        "metadata_quote_type": quote_type,
    }


def resolve_holdings(df):
    records = []

    for _, row in df.iterrows():
        resolved = resolve_single_holding(row["identifier"], row["description"])
        merged = row.to_dict()
        merged.update(resolved)
        records.append(merged)

    return pd.DataFrame(records)


def calculate_asset_allocation(df):
    total = df["gbp_value"].sum()
    alloc = df.groupby("asset_class", dropna=False)["gbp_value"].sum().reset_index()
    alloc["weight_pct"] = (alloc["gbp_value"] / total) * 100 if total else 0
    alloc["asset_class"] = pd.Categorical(alloc["asset_class"], ASSET_CLASS_ORDER, ordered=True)
    alloc = alloc.sort_values(["asset_class", "gbp_value"], ascending=[True, False]).reset_index(drop=True)
    alloc["asset_class"] = alloc["asset_class"].astype(str)
    return alloc


def calculate_region_allocation(df):
    total = df["gbp_value"].sum()
    alloc = df.groupby("region", dropna=False)["gbp_value"].sum().reset_index()
    alloc["weight_pct"] = (alloc["gbp_value"] / total) * 100 if total else 0
    alloc["region"] = pd.Categorical(alloc["region"], REGION_ORDER, ordered=True)
    alloc = alloc.sort_values(["region", "gbp_value"], ascending=[True, False]).reset_index(drop=True)
    alloc["region"] = alloc["region"].astype(str)
    return alloc


def calculate_portfolio_trailing_returns(df):
    matched = df[df["matched_ticker"].notna()].copy()

    if matched.empty:
        return pd.DataFrame(
            {
                "period": ["1Y", "3Y", "5Y"],
                "portfolio_return_pct": [None, None, None],
            }
        )

    total_value = matched["gbp_value"].sum()
    matched["weight"] = matched["gbp_value"] / total_value if total_value else 0

    results = []

    for years, label in [(1, "1Y"), (3, "3Y"), (5, "5Y")]:
        weighted_return = 0.0
        valid_weight = 0.0

        for _, row in matched.iterrows():
            series = get_price_history(row["matched_ticker"], years)

            if series.empty or len(series) < 2:
                continue

            holding_return = (series.iloc[-1] / series.iloc[0]) - 1
            weighted_return += holding_return * row["weight"]
            valid_weight += row["weight"]

        portfolio_return = (weighted_return / valid_weight) if valid_weight > 0 else None

        results.append(
            {
                "period": label,
                "portfolio_return_pct": round(portfolio_return * 100, 2) if portfolio_return is not None else None,
            }
        )

    return pd.DataFrame(results)


def create_portfolio_pptx(client_name, returns_df):
    prs = Presentation()

    # Slide 1: Title slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = f"Portfolio Review for {client_name}"
    slide1.placeholders[1].text = f"Generated on {date.today().strftime('%d %B %Y')}"

    # Slide 2: Performance table
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # Title only
    slide2.shapes.title.text = "Portfolio Performance"

    rows = len(returns_df) + 1
    cols = 2
    left = Inches(1.2)
    top = Inches(1.8)
    width = Inches(7.0)
    height = Inches(2.0)

    table = slide2.shapes.add_table(rows, cols, left, top, width, height).table

    table.columns[0].width = Inches(3.0)
    table.columns[1].width = Inches(4.0)

    table.cell(0, 0).text = "Period"
    table.cell(0, 1).text = "Portfolio Return (%)"

    for i, (_, row) in enumerate(returns_df.iterrows(), start=1):
        table.cell(i, 0).text = str(row["period"])
        if pd.notna(row["portfolio_return_pct"]):
            table.cell(i, 1).text = f"{row['portfolio_return_pct']:.2f}%"
        else:
            table.cell(i, 1).text = "N/A"

    for row in table.rows:
        for cell in row.cells:
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                for run in paragraph.runs:
                    run.font.size = Pt(14)

    # Slide 3: Thank you
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    slide3.shapes.title.text = "Thank you for your time"

    return prs


st.title("Investment Portfolio Analysis Tool")
st.write(
    "Upload a CSV of your holdings, confirm the extraction, then generate allocation, return, and geographic insights using Yahoo Finance metadata where available."
)

with st.expander("Expected CSV format"):
    st.markdown(
        """
        Your file should contain these columns or close equivalents:
        - Ticker / ISIN / SEDOL
        - Description
        - GBP Value

        Example:
        ```csv
        Ticker,Description,GBP Value
        PYPL,PayPal Holdings Inc,5000
        CROX,Crocs Inc,3000
        ```
        """
    )

uploaded_file = st.file_uploader("Upload holdings CSV", type=["csv"])

if "parsed_df" not in st.session_state:
    st.session_state.parsed_df = None
if "confirmed" not in st.session_state:
    st.session_state.confirmed = False

if uploaded_file is not None:
    try:
        parsed_df = parse_uploaded_holdings(uploaded_file)
        st.session_state.parsed_df = parsed_df

        st.subheader("Step 1: Confirm extracted holdings")
        st.dataframe(parsed_df, use_container_width=True)

        if st.button("Confirm extraction"):
            st.session_state.confirmed = True

    except Exception as e:
        st.error(f"Could not parse file: {e}")

if st.session_state.confirmed and st.session_state.parsed_df is not None:
    with st.spinner("Resolving holdings and fetching Yahoo Finance metadata..."):
        enriched_df = resolve_holdings(st.session_state.parsed_df)

    st.subheader("Step 2: Enriched holdings")
    st.dataframe(enriched_df, use_container_width=True)

    unresolved = enriched_df[
        (enriched_df["region"] == "Unclassified") | (enriched_df["asset_class"] == "Unclassified")
    ]
    if not unresolved.empty:
        st.warning(
            f"{len(unresolved)} holding(s) could not be fully classified from Yahoo Finance metadata. "
            "They are still included, but some values may remain Unclassified."
        )

    asset_alloc_df = calculate_asset_allocation(enriched_df)
    region_alloc_df = calculate_region_allocation(enriched_df)
    returns_df = calculate_portfolio_trailing_returns(enriched_df)

    st.subheader("Step 3: Portfolio insights")

    col1, col2 = st.columns(2)

    with col1:
        st.markdown("### Asset allocation")
        fig_asset = px.pie(
            asset_alloc_df,
            names="asset_class",
            values="gbp_value",
            hole=0.45,
        )
        st.plotly_chart(fig_asset, use_container_width=True)
        st.dataframe(asset_alloc_df, use_container_width=True)

    with col2:
        st.markdown("### Geographic allocation")
        fig_region = px.bar(
            region_alloc_df,
            x="region",
            y="gbp_value",
            text="weight_pct",
        )
        fig_region.update_traces(texttemplate="%{text:.2f}%")
        st.plotly_chart(fig_region, use_container_width=True)
        st.dataframe(region_alloc_df, use_container_width=True)

    st.markdown("### Trailing returns")
    st.dataframe(returns_df, use_container_width=True)

    st.subheader("Step 4: Generate PowerPoint")
    client_name = st.text_input("Enter client name")

    if st.button("Create portfolio review PPTX"):
        if not client_name.strip():
            st.warning("Please enter the client name first.")
        else:
            with st.spinner("Generating PowerPoint presentation..."):
                prs = create_portfolio_pptx(client_name.strip(), returns_df)
                pptx_buffer = io.BytesIO()
                prs.save(pptx_buffer)
                pptx_buffer.seek(0)

            safe_name = client_name.strip().replace(" ", "_")
            st.success("PowerPoint presentation created successfully.")
            st.download_button(
                label="Download portfolio review PPTX",
                data=pptx_buffer.getvalue(),
                file_name=f"portfolio_review_{safe_name}.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )
